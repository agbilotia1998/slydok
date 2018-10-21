from pptx import Presentation
# from googleapiclient import http
# from googleapiclient.discovery import build
# from httplib2 import Http
# from oauth2client import file, client, tools
from io import BytesIO, FileIO
from docx import Document
from pptx.util import Pt
from conversion import gensim_summarizer
import docx
from pptx.util import Inches
from images import extract_images
import scipy.ndimage
import re
from docx2csv import extract_tables, extract

def select_ppt(selection):
    switcher = {
        1: 'templates/pitch.pptx',
        2: 'templates/your_big_idea.pptx',
        3: 'templates/consulting_proposal.pptx',
        4: 'templates/recipe_book.pptx',
        5: 'templates/portfolio.pptx'
    }
    return switcher.get(selection)

filename = input("Enter path to the file that you want to convert: ")
document = Document(filename)

selection = input(
    'Select among the following themes:\n\t1. Pitch\n\t2. Your big idea\n\t3. Recipe Book\n\t4. Consulting Proposal\n\t5. Portfolio\n\nPlease select by entering the number: ')

ppt_name = select_ppt(selection)

prs = Presentation(ppt_name)
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
body_shape = None
tf = None
images = extract_images(filename)
image_count = 1
default_constraint = 900
constraint = 900

def add_table(table_position, table_list):
    title_only_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(title_only_slide_layout)
    shapes = slide.shapes
    print table_position, table_list[0]
    rows = len(table_list[0])
    cols = len(table_list[0][0])
    left = top = Inches(2.0)
    width = Inches(6.0)
    height = Inches(2.8)

    table = shapes.add_table(rows, cols, left, top, width, height).table

    for indr, rows in enumerate(table_list[0]):
        for indc, col in enumerate(rows):
                table.cell(indr, indc).text = col
    table_list.pop(0)

table_indexes = []

def get_table_indexes():
    for ind, el in enumerate(document.element.body):
        if 'tblGrid' in dir(el):
            table_indexes.append(ind)

get_table_indexes()

def add_overflow_slide(result):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    body_shape = slide.shapes.add_textbox(Inches(0.5),Inches(1.15),Inches(8.8),Inches(0.0))
    global tf
    tf = body_shape.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = result
    p.level = 0
    global constraint
    constraint = default_constraint - len(result)
    print constraint

def emu_to_pixels(emu):
    return int(round(emu / 9525.0))

def pixels_to_emu(pixels):
    return int(pixels * 9525)

def iter_headings(paragraphs):
    for paragraph in paragraphs:
        if paragraph.style.name.startswith('Heading'):
            yield paragraph

my_list = list()
ct = 0
head = ''
present = 0

tables = extract_tables(filename)
print tables

for ind, para in enumerate(document.paragraphs):
    if ind in table_indexes:
        add_table(ind, tables)
        continue

    my_list.append(para.style.name)

    if para.style.name == 'Title':
        title.text = para.text
    if para.style.name == 'Subtitle':
        subtitle.text = para.text
    if para.style.name.startswith('Heading'):
        if(not my_list[ct-1].startswith('Heading')):
            layout = prs.slide_layouts[5]
            slide = prs.slides.add_slide(layout)
            # slide.placeholders[1].getparent().remove(slide.placeholders[1])
            heading = slide.shapes.title
            heading.text = para.text
            constraint = default_constraint

            body_shape = slide.shapes.add_textbox(Inches(0.5),Inches(1.15),Inches(8.8),Inches(0.0))
            tf = body_shape.text_frame
            tf.word_wrap = True
        else:
            p = tf.add_paragraph()
            p.font.bold = True
            p.font.size = Pt(16)
            #head += para.text + '\n'
            p.text = para.text

    if(para.style.name == 'Normal'):
        result = head + para.text
        if len(result) > constraint:
            head = ''
            p = tf.add_paragraph()
            p.font.bold = True
            p.text = result[:constraint]
            p.level = 0
            add_overflow_slide(result[constraint + 1:])
        else:
            constraint -= len(result)
            head = ''
            p = tf.add_paragraph()
            p.text = result
            p.level = 0

    if len(para._p.r_lst[0].drawing_lst):
        image_name = 'image{0}.jpeg'.format(image_count)
        image_count += 1
        height, width, channels = scipy.ndimage.imread(image_name).shape

        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        left = (prs.slide_width - pixels_to_emu(width)) / 3
        top = (prs.slide_height - pixels_to_emu(height)) / 3
        pic = slide.shapes.add_picture(image_name, left, top)

    ct+=1
    
prs.save('test.pptx')

print('Generated Your PPTX file succesfully.')